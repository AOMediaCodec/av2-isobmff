"""Meeting contribution slide generator.

Generates a JVET/MPEG/AOM-style .pptx presentation from a built spec or
from a subset of sections, suitable for standards meeting contributions.
"""

from __future__ import annotations

import logging
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from pptx import Presentation as PptxPresentation


# ---------------------------------------------------------------------------
# Data model
# ---------------------------------------------------------------------------


@dataclass
class SlideContent:
    section_id: str
    heading: str
    bullets: list[str] = field(default_factory=list)
    notes: str = ""


# ---------------------------------------------------------------------------
# Content extraction
# ---------------------------------------------------------------------------


def _sentences(text: str, max_count: int = 4) -> list[str]:
    """Split *text* into at most *max_count* sentence fragments."""
    parts = re.split(r"(?<=[.!?])\s+", text.strip())
    return [p.strip() for p in parts[:max_count] if p.strip()]


def extract_slides_content(
    soup,
    section_ids: list[str] | None = None,
) -> list[SlideContent]:
    """Extract slide-ready content from *soup*.

    Args:
        soup:        BeautifulSoup of the built spec HTML.
        section_ids: Explicit list of section ids to include.  If None,
                     all top-level ``<section>`` elements are used.

    Returns:
        List of :class:`SlideContent` objects in document order.
    """
    slides: list[SlideContent] = []

    def _process_section(elem) -> SlideContent | None:
        sec_id = elem.get("id", "")
        heading_el = elem.find(re.compile(r"^h[1-6]$"))
        if not heading_el:
            return None
        heading = heading_el.get_text(strip=True)[:120]

        # Collect paragraph text from direct children (not sub-sections)
        body_text = ""
        for child in elem.children:
            child_name = getattr(child, "name", None)
            if child_name in (None, "section"):
                continue
            if child_name and re.match(r"^h[1-6]$", child_name):
                continue
            body_text += child.get_text(" ", strip=True) + " "

        bullets = _sentences(body_text.strip(), max_count=5)
        return SlideContent(section_id=sec_id, heading=heading, bullets=bullets)

    if section_ids:
        for sid in section_ids:
            elem = soup.find(id=sid)
            if elem:
                sc = _process_section(elem)
                if sc:
                    slides.append(sc)
    else:
        # Top-level sections only
        for section in soup.find_all("section", recursive=False):
            sc = _process_section(section)
            if sc:
                slides.append(sc)
        if not slides:
            # Fallback: h2-level headings
            for h2 in soup.find_all("h2"):
                parent = h2.parent
                if parent:
                    sc = _process_section(parent)
                    if sc:
                        slides.append(sc)

    logging.info(f"Extracted {len(slides)} slide(s) from spec")
    return slides


# ---------------------------------------------------------------------------
# Presentation building
# ---------------------------------------------------------------------------


def _pt(points: float):
    """Convert points to EMU (English Metric Units used by python-pptx)."""
    from pptx.util import Pt

    return Pt(points)


def _emu(cm: float):
    from pptx.util import Cm

    return Cm(cm)


def _rgb(r: int, g: int, b: int):
    from pptx.dml.color import RGBColor

    return RGBColor(r, g, b)


# AOM/JVET-inspired palette
_COLOR_HEADER_BG = (3, 69, 117)  # dark blue
_COLOR_HEADER_FG = (255, 255, 255)  # white
_COLOR_BODY_FG = (26, 26, 26)  # near-black
_COLOR_ACCENT = (0, 102, 204)  # link blue
_COLOR_SLIDE_BG = (255, 255, 255)  # white


def _add_title_slide(prs, title: str, subtitle: str, date: str) -> None:
    """Add the opening title slide."""
    from pptx.util import Pt

    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    # Header bar
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        _emu(0),
        _emu(0),
        prs.slide_width,
        _emu(4),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(*_COLOR_HEADER_BG)
    bar.line.fill.background()

    # Title text
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = _rgb(*_COLOR_HEADER_FG)

    # Subtitle
    txb = slide.shapes.add_textbox(_emu(1), _emu(4.5), prs.slide_width - _emu(2), _emu(2))
    tf2 = txb.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = _rgb(*_COLOR_BODY_FG)

    # Date (bottom right)
    date_box = slide.shapes.add_textbox(
        prs.slide_width - _emu(6),
        prs.slide_height - _emu(1.5),
        _emu(5.5),
        _emu(1),
    )
    date_tf = date_box.text_frame
    date_tf.paragraphs[0].text = date
    date_tf.paragraphs[0].font.size = Pt(12)
    date_tf.paragraphs[0].font.color.rgb = _rgb(120, 120, 120)


def _add_agenda_slide(prs, slides_content: list[SlideContent]) -> None:
    """Add an agenda slide listing all sections."""
    from pptx.util import Pt

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Header
    bar = slide.shapes.add_shape(1, _emu(0), _emu(0), prs.slide_width, _emu(1.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(*_COLOR_HEADER_BG)
    bar.line.fill.background()
    tf = bar.text_frame
    tf.paragraphs[0].text = "Agenda"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = _rgb(*_COLOR_HEADER_FG)

    # Bullet list
    txb = slide.shapes.add_textbox(_emu(1), _emu(2.2), prs.slide_width - _emu(2), _emu(12))
    tf2 = txb.text_frame
    tf2.word_wrap = True
    for i, sc in enumerate(slides_content):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"• {sc.heading}"
        p.font.size = Pt(16)
        p.font.color.rgb = _rgb(*_COLOR_BODY_FG)


def _add_content_slide(prs, sc: SlideContent) -> None:
    """Add a content slide for one section."""
    from pptx.util import Pt

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Header bar
    bar = slide.shapes.add_shape(1, _emu(0), _emu(0), prs.slide_width, _emu(1.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(*_COLOR_HEADER_BG)
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = sc.heading
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = _rgb(*_COLOR_HEADER_FG)

    # Content bullets
    txb = slide.shapes.add_textbox(_emu(1), _emu(2.2), prs.slide_width - _emu(2), _emu(12))
    tf2 = txb.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(sc.bullets):
        p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p2.text = f"• {bullet}"
        p2.font.size = Pt(15)
        p2.font.color.rgb = _rgb(*_COLOR_BODY_FG)

    if not sc.bullets:
        p2 = tf2.paragraphs[0]
        p2.text = "(see specification)"
        p2.font.size = Pt(15)
        p2.font.color.rgb = _rgb(120, 120, 120)
        p2.font.italic = True

    # Section ID footer
    footer = slide.shapes.add_textbox(
        _emu(0.5),
        prs.slide_height - _emu(1),
        _emu(6),
        _emu(0.8),
    )
    footer.text_frame.paragraphs[0].text = f"§ {sc.section_id}"
    footer.text_frame.paragraphs[0].font.size = Pt(9)
    footer.text_frame.paragraphs[0].font.color.rgb = _rgb(160, 160, 160)


def build_presentation(
    slides_content: list[SlideContent],
    metadata: dict[str, str],
) -> PptxPresentation:
    """Build a python-pptx Presentation from *slides_content*.

    Args:
        slides_content: From :func:`extract_slides_content`.
        metadata:       Dict with keys: ``title``, ``subtitle``, ``date``.

    Returns:
        A ``pptx.Presentation`` object ready to save.
    """
    try:
        from pptx import Presentation
    except ImportError:
        logging.error("python-pptx is required: pip install python-pptx")
        raise SystemExit(1)

    prs = Presentation()
    prs.slide_width = _emu(33.87)  # 16:9 widescreen (33.87 × 19.05 cm)
    prs.slide_height = _emu(19.05)

    title = metadata.get("title", "Specification")
    subtitle = metadata.get("subtitle", "")
    date = metadata.get("date", "")

    _add_title_slide(prs, title, subtitle, date)
    if slides_content:
        _add_agenda_slide(prs, slides_content)
    for sc in slides_content:
        _add_content_slide(prs, sc)

    return prs


def save_presentation(prs: PptxPresentation, path: Path) -> None:
    """Save *prs* to *path*."""
    try:
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(path))
        logging.info(f"Slides saved: {path}")
    except OSError as exc:
        logging.error(f"Failed to save presentation to {path}: {exc}")
