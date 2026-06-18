# \!/usr/bin/env python3
"""Generate the specbuild feature slide deck.

Creates a complete 15-slide PowerPoint deck documenting all specbuild features
with editable vector graphics. Widescreen 16:9, white background.

Usage:
    python scripts/generate_slides.py [output.pptx]
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ── Theme colors ─────────────────────────────────────────────────────────────

DARK = RGBColor(0x1A, 0x20, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x2D, 0x37, 0x48)
TEXT = RGBColor(0x33, 0x33, 0x33)
TEXT_LIGHT = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0x29, 0x6B, 0xCF)  # blue
ACCENT2 = RGBColor(0x2F, 0x9E, 0x5A)  # green
ACCENT3 = RGBColor(0xD9, 0x7B, 0x0B)  # amber
ACCENT4 = RGBColor(0x7C, 0x3A, 0xED)  # purple
ACCENT5 = RGBColor(0xDB, 0x44, 0x37)  # red
TEAL = RGBColor(0x06, 0x94, 0xA2)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFA)
CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
MUTED = RGBColor(0x71, 0x80, 0x96)
CODE_BG = RGBColor(0xF1, 0xF5, 0xF9)

# Change-bar / keyword / report colors
CB_GREEN = RGBColor(0xC6, 0xF6, 0xD5)
CB_GREEN_DARK = RGBColor(0x27, 0x67, 0x49)
CB_YELLOW = RGBColor(0xFE, 0xFC, 0xBF)
CB_YELLOW_DARK = RGBColor(0x97, 0x5A, 0x16)
KW_SHALL = RGBColor(0xC5, 0x3D, 0x30)
KW_SHOULD = RGBColor(0xDD, 0x6B, 0x20)
KW_MAY = RGBColor(0x29, 0x6B, 0xCF)
REPORT_OK = RGBColor(0x27, 0x67, 0x49)
REPORT_WARN = RGBColor(0x97, 0x5A, 0x16)

# ── Slide dimensions (widescreen 16:9) ──────────────────────────────────────

SLD_W = Inches(13.333)
SLD_H = Inches(7.5)


# ═════════════════════════════════════════════════════════════════════════════
# Helpers
# ═════════════════════════════════════════════════════════════════════════════


def add_shape(
    slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None
):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=14,
    bold=False,
    color=TEXT,
    alignment=PP_ALIGN.LEFT,
    font_name="Calibri",
    italic=False,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_text(slide, left, top, width, height, runs, font_size=12, font_name="Calibri"):
    """Add text box with multiple styled runs.
    runs: list of (text, color, bold, italic, font_name_override)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    for i, run_data in enumerate(runs):
        text = run_data[0]
        color = run_data[1] if len(run_data) > 1 else TEXT
        bold = run_data[2] if len(run_data) > 2 else False
        italic_ = run_data[3] if len(run_data) > 3 else False
        fn = run_data[4] if len(run_data) > 4 else font_name
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic_
        run.font.name = fn
    return txBox


def add_bullet_text(
    slide,
    left,
    top,
    width,
    height,
    items,
    font_size=13,
    color=TEXT,
    spacing=Pt(4),
    font_name="Calibri",
    bullet_char="\u2022",
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn("a:buChar"), {"char": bullet_char})
        for child in pPr.findall(qn("a:buChar")):
            pPr.remove(child)
        for child in pPr.findall(qn("a:buNone")):
            pPr.remove(child)
        pPr.append(buChar)
    return txBox


def add_card(
    slide,
    left,
    top,
    width,
    height,
    title,
    body_items,
    accent_color=ACCENT,
    title_size=14,
    body_size=11,
):
    """Rounded-rect card with accent top bar."""
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        fill_color=CARD_BG,
        line_color=BORDER,
        line_width=Pt(1),
    )
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, Pt(4), fill_color=accent_color)
    add_text_box(
        slide,
        left + Inches(0.15),
        top + Inches(0.1),
        width - Inches(0.3),
        Inches(0.35),
        title,
        font_size=title_size,
        bold=True,
        color=accent_color,
    )
    if body_items:
        add_bullet_text(
            slide,
            left + Inches(0.15),
            top + Inches(0.45),
            width - Inches(0.3),
            height - Inches(0.55),
            body_items,
            font_size=body_size,
            color=TEXT_LIGHT,
        )


def add_icon_circle(slide, left, top, size, color, label, label_size=20):
    add_shape(slide, MSO_SHAPE.OVAL, left, top, size, size, fill_color=color)
    txBox = slide.shapes.add_textbox(left, top, size, size)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(label_size)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(label_size * 0.6)


def slide_title_bar(slide, title, subtitle=None):
    """Consistent title area at top of content slides."""
    add_shape(
        slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLD_W, Inches(1.15), fill_color=NEAR_BLACK
    )
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.15), SLD_W, Pt(4), fill_color=ACCENT)
    add_text_box(
        slide,
        Inches(0.6),
        Inches(0.15),
        Inches(10),
        Inches(0.5),
        title,
        font_size=28,
        bold=True,
        color=WHITE,
    )
    if subtitle:
        add_text_box(
            slide,
            Inches(0.6),
            Inches(0.62),
            Inches(11),
            Inches(0.4),
            subtitle,
            font_size=14,
            color=RGBColor(0xA0, 0xAE, 0xC0),
        )


def mockup_frame(slide, left, top, width, height, title, frame_color=BORDER):
    """Draw a browser/app-like frame around a mockup area."""
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        fill_color=WHITE,
        line_color=frame_color,
        line_width=Pt(1.5),
    )
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        left,
        top,
        width,
        Inches(0.35),
        fill_color=LIGHT_BG,
        line_color=frame_color,
        line_width=Pt(1),
    )
    for i, c in enumerate([ACCENT5, ACCENT3, ACCENT2]):
        add_shape(
            slide,
            MSO_SHAPE.OVAL,
            left + Inches(0.12) + i * Inches(0.22),
            top + Inches(0.1),
            Inches(0.14),
            Inches(0.14),
            fill_color=c,
        )
    add_text_box(
        slide,
        left + Inches(0.8),
        top + Inches(0.04),
        width - Inches(1.0),
        Inches(0.3),
        title,
        font_size=10,
        bold=True,
        color=MUTED,
        alignment=PP_ALIGN.CENTER,
    )
    return left, top + Inches(0.35)


def add_label(slide, left, top, width, text, color=ACCENT, font_size=11):
    """Add a small label/badge."""
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.28), fill_color=color)
    add_text_box(
        slide,
        left,
        top + Inches(0.02),
        width,
        Inches(0.24),
        text,
        font_size=font_size,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )


# ═════════════════════════════════════════════════════════════════════════════
# Slide builders — each function creates one slide
# ═════════════════════════════════════════════════════════════════════════════


def slide_01_title(prs, layout):
    """Title slide."""
    slide = prs.slides.add_slide(layout)

    add_shape(
        slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), SLD_H, fill_color=ACCENT
    )

    add_text_box(
        slide,
        Inches(0.8),
        Inches(1.2),
        Inches(8),
        Inches(1.0),
        "specbuild",
        font_size=54,
        bold=True,
        color=NEAR_BLACK,
    )
    add_text_box(
        slide,
        Inches(0.8),
        Inches(2.2),
        Inches(10),
        Inches(0.8),
        "A Reusable Pipeline for Building Bikeshed Specifications",
        font_size=24,
        color=ACCENT,
    )
    add_text_box(
        slide,
        Inches(0.8),
        Inches(3.1),
        Inches(10),
        Inches(0.5),
        "From source to publication-ready HTML, PDF, Word, and LaTeX",
        font_size=16,
        color=MUTED,
    )

    badges = [
        ("13", "Quality\nChecks", ACCENT),
        ("18", "Enhancements", ACCENT2),
        ("20", "Output\nTasks", ACCENT3),
        ("6", "Build\nProfiles", ACCENT4),
    ]
    badge_left = Inches(0.8)
    for count, label, color in badges:
        add_icon_circle(slide, badge_left, Inches(4.3), Inches(1.15), color, count, label_size=30)
        add_text_box(
            slide,
            badge_left - Inches(0.15),
            Inches(5.55),
            Inches(1.45),
            Inches(0.6),
            label,
            font_size=10,
            color=MUTED,
            alignment=PP_ALIGN.CENTER,
        )
        badge_left += Inches(1.7)

    add_text_box(
        slide,
        Inches(0.8),
        Inches(6.6),
        Inches(6),
        Inches(0.4),
        "v0.1.0   \u2022   March 2026   \u2022   Python 3.10+   \u2022   BSD-3-Clause",
        font_size=12,
        color=MUTED,
    )

    # Pipeline phases panel
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.0),
        Inches(1.5),
        Inches(3.8),
        Inches(5.0),
        fill_color=LIGHT_BG,
        line_color=BORDER,
        line_width=Pt(1),
    )
    add_text_box(
        slide,
        Inches(9.2),
        Inches(1.7),
        Inches(3.4),
        Inches(0.4),
        "PIPELINE PHASES",
        font_size=12,
        bold=True,
        color=ACCENT,
    )

    phase_labels = [
        ("\u2776  Merge .bs sources", ACCENT),
        ("\u2777  Compile with Bikeshed", ACCENT),
        ("\u2778  Run quality checks", ACCENT2),
        ("\u2779  Apply enhancements", ACCENT3),
        ("\u277a  Generate outputs", ACCENT4),
    ]
    for i, (label, color) in enumerate(phase_labels):
        y = Inches(2.2) + i * Inches(0.75)
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(9.3),
            y,
            Inches(3.4),
            Inches(0.55),
            fill_color=WHITE,
            line_color=color,
            line_width=Pt(2),
        )
        add_text_box(
            slide,
            Inches(9.5),
            y + Inches(0.1),
            Inches(3.0),
            Inches(0.35),
            label,
            font_size=13,
            bold=True,
            color=color,
        )


def slide_02_architecture(prs, layout):
    """Architecture Overview."""
    slide = prs.slides.add_slide(layout)
    slide_title_bar(
        slide,
        "Architecture Overview",
        "Plugin-based pipeline with five phases and concurrent execution",
    )

    phases = [
        ("1  MERGE", "Combine .bs files\ninto index.bs\nvia manifest order", ACCENT),
        ("2  COMPILE", "Bikeshed \u2192 HTML\nSDL table conversion\nVersion stamping", ACCENT),
        ("3  CHECK", "13 quality checks\nConcurrent (read-only)\nThreadPoolExecutor", ACCENT2),
        ("4  ENHANCE", "18 enhancements\nSequential mutations\nOrdered by priority", ACCENT3),
        ("5  OUTPUT", "20 output tasks\nParallelizable\nPDF, Word, LaTeX, \u2026", ACCENT4),
    ]

    box_w = Inches(2.15)
    box_h = Inches(1.8)
    start_x = Inches(0.4)
    y = Inches(1.7)

    for i, (title, desc, color) in enumerate(phases):
        x = start_x + i * Inches(2.5)
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            y,
            box_w,
            box_h,
            fill_color=CARD_BG,
            line_color=color,
            line_width=Pt(2),
        )
        add_text_box(
            slide,
            x + Inches(0.12),
            y + Inches(0.12),
            box_w - Inches(0.24),
            Inches(0.35),
            title,
            font_size=14,
            bold=True,
            color=color,
        )
        add_text_box(
            slide,
            x + Inches(0.12),
            y + Inches(0.55),
            box_w - Inches(0.24),
            Inches(1.1),
            desc,
            font_size=11,
            color=TEXT_LIGHT,
        )
        if i < len(phases) - 1:
            arrow_x = x + box_w + Inches(0.03)
            add_shape(
                slide,
                MSO_SHAPE.RIGHT_ARROW,
                arrow_x,
                y + Inches(0.55),
                Inches(0.3),
                Inches(0.5),
                fill_color=BORDER,
            )

    points = [
        "Quality checks share a single BeautifulSoup parse \u2014 concurrent via ThreadPoolExecutor",
        "Enhancements run sequentially with ordered priorities (BeautifulSoup not thread-safe for writes)",
        "Output tasks are independent \u2014 enable --parallel-outputs for concurrent execution",
        "Plugin registry: 51+ plugins discovered via decorators, listed with --help-features",
        "TOML configuration: specbuild.toml or [tool.specbuild] in pyproject.toml",
    ]
    add_text_box(
        slide,
        Inches(0.6),
        Inches(3.85),
        Inches(5),
        Inches(0.3),
        "KEY DESIGN PRINCIPLES",
        font_size=13,
        bold=True,
        color=ACCENT,
    )
    add_bullet_text(
        slide,
        Inches(0.6),
        Inches(4.2),
        Inches(12),
        Inches(2.8),
        points,
        font_size=12,
        color=TEXT_LIGHT,
        spacing=Pt(5),
    )


def slide_03_quality_checks(prs, layout):
    """Quality Checks (13)."""
    slide = prs.slides.add_slide(layout)
    slide_title_bar(
        slide,
        "Quality Checks (13)",
        "Concurrent read-only analysis  \u2022  Each supports --strict mode  \u2022  Enable all: --all-checks",
    )

    checks = [
        ("Cross-References", ["--validate-refs", "Broken internal links"], ACCENT),
        ("SDL References", ["--validate-sdl-refs", "Unresolved function calls"], ACCENT),
        ("Image Check", ["--check-images", "Missing/broken images"], ACCENT),
        ("Table Validation", ["--check-tables", "Columns, headers, scope"], ACCENT),
        ("Referenceable", ["--check-referenceable", "Caption IDs on tables/figures"], ACCENT2),
        ("Accessibility", ["--accessibility-audit", "WCAG compliance audit"], ACCENT2),
        ("Definitions", ["--check-dfn", "Unused/undefined <dfn>"], ACCENT2),
        ("Terminology", ["--check-terminology", "Inconsistent term usage"], ACCENT2),
        ("Orphan Refs", ["--check-orphan-refs", "Uncited bibliography entries"], ACCENT3),
        ("External Links", ["--check-links", "URL reachability check"], ACCENT3),
        ("Spell Check", ["--spellcheck", "Domain-aware spelling"], ACCENT3),
        ("Duplicates", ["--check-duplicates", "Near-duplicate paragraphs"], ACCENT3),
        ("Editorial", ["--editorial", "Compound words, caps"], ACCENT4),
    ]

    cols = 4
    card_w = Inches(2.95)
    card_h = Inches(1.12)
    start_x = Inches(0.45)
    start_y = Inches(1.55)

    for i, (title, items, color) in enumerate(checks):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + Inches(0.15))
        y = start_y + row * (card_h + Inches(0.12))
        add_card(
            slide,
            x,
            y,
            card_w,
            card_h,
            title,
            items,
            accent_color=color,
            title_size=12,
            body_size=10,
        )

    add_text_box(
        slide,
        Inches(0.5),
        Inches(6.2),
        Inches(11),
        Inches(0.5),
        "Strict mode (--<name>-strict) exits with error code 1 on findings \u2014 ideal for CI gates",
        font_size=12,
        color=MUTED,
    )


def slide_04_qc_details(prs, layout):
    """Quality Check Reports: Debugging Aids."""
    slide = prs.slides.add_slide(layout)
    slide_title_bar(
        slide,
        "Quality Check Reports: Debugging Aids",
        "Detailed findings with context and suggestions for each quality check",
    )

    # --- Table Validation ---
    tv_x, tv_y = Inches(0.4), Inches(1.5)
    tv_w, tv_h = Inches(6.2), Inches(2.6)
    _, tv_cy = mockup_frame(
        slide, tv_x, tv_y, tv_w, tv_h, "Table Validation  \u2014  --check-tables"
    )

    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        tv_x + Inches(0.1),
        tv_cy + Inches(0.1),
        Inches(6.0),
        Inches(0.28),
        fill_color=NEAR_BLACK,
    )
    for cx2, ct in [
        (tv_x + Inches(0.15), "Rule"),
        (tv_x + Inches(1.4), "Detail"),
        (tv_x + Inches(4.0), "Near Section"),
    ]:
        add_text_box(
            slide,
            cx2,
            tv_cy + Inches(0.12),
            Inches(2.5),
            Inches(0.24),
            ct,
            font_size=9,
            bold=True,
            color=WHITE,
        )

    table_issues = [
        ("table-ragged-row", "Row 5 has 3 columns, expected 4", "\u00a74.2 Operators"),
        ("th-no-scope", "<th> without scope: Mode", "\u00a75.1 Syntax"),
        ("table-no-caption", "Table has no <caption> or aria-label", "\u00a76.3 Loop Filter"),
        ("table-no-thead", "Has <th> but no <thead>/<tbody>", "\u00a76.8 Transform"),
    ]
    try3 = tv_cy + Inches(0.4)
    for i, (rule, detail, section) in enumerate(table_issues):
        bg = CARD_BG if i % 2 == 0 else WHITE
        add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            tv_x + Inches(0.1),
            try3,
            Inches(6.0),
            Inches(0.35),
            fill_color=bg,
        )
        rule_color = ACCENT5 if "ragged" in rule else (ACCENT3 if "no-" in rule else ACCENT)
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            tv_x + Inches(0.15),
            try3 + Inches(0.05),
            Inches(1.15),
            Inches(0.24),
            fill_color=rule_color,
        )
        add_text_box(
            slide,
            tv_x + Inches(0.15),
            try3 + Inches(0.05),
            Inches(1.15),
            Inches(0.24),
            rule,
            font_size=9,
            bold=True,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
            font_name="Consolas",
        )
        add_text_box(
            slide,
            tv_x + Inches(1.4),
            try3 + Inches(0.05),
            Inches(2.5),
            Inches(0.24),
            detail,
            font_size=9,
            color=TEXT,
        )
        add_text_box(
            slide,
            tv_x + Inches(4.0),
            try3 + Inches(0.05),
            Inches(2.0),
            Inches(0.24),
            section,
            font_size=9,
            color=TEXT_LIGHT,
        )
        try3 += Inches(0.35)

    # --- Terminology Check ---
    term_x, term_y = Inches(6.8), Inches(1.5)
    term_w, term_h = Inches(6.1), Inches(2.6)
    _, term_cy = mockup_frame(
        slide, term_x, term_y, term_w, term_h, "Terminology Check  \u2014  --check-terminology"
    )

    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        term_x + Inches(0.1),
        term_cy + Inches(0.1),
        Inches(5.9),
        Inches(0.28),
        fill_color=NEAR_BLACK,
    )
    for cx2, ct, cw2 in [
        (term_x + Inches(0.15), "Canonical Term", Inches(1.6)),
        (term_x + Inches(1.8), "Variant Found", Inches(1.6)),
        (term_x + Inches(3.5), "Canonical #", Inches(0.7)),
        (term_x + Inches(4.3), "Variant #", Inches(0.7)),
        (term_x + Inches(5.1), "Action", Inches(0.8)),
    ]:
        add_text_box(
            slide,
            cx2,
            term_cy + Inches(0.12),
            cw2,
            Inches(0.24),
            ct,
            font_size=9,
            bold=True,
            color=WHITE,
        )

    term_issues = [
        ("intra frame", "intra-frame", "142", "3", "Standardize"),
        ("quantization", "quantisation", "89", "7", "US spelling"),
        ("color space", "colour space", "45", "12", "US spelling"),
        ("bitstream", "bit stream", "234", "5", "Use compound"),
    ]
    try4 = term_cy + Inches(0.4)
    for i, (canonical, variant, c_count, v_count, action) in enumerate(term_issues):
        bg = CARD_BG if i % 2 == 0 else WHITE
        add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            term_x + Inches(0.1),
            try4,
            Inches(5.9),
            Inches(0.35),
            fill_color=bg,
        )
        add_text_box(
            slide,
            term_x + Inches(0.15),
            try4 + Inches(0.05),
            Inches(1.5),
            Inches(0.24),
            canonical,
            font_size=9,
            bold=True,
            color=ACCENT2,
        )
        add_text_box(
            slide,
            term_x + Inches(1.8),
            try4 + Inches(0.05),
            Inches(1.5),
            Inches(0.24),
            variant,
            font_size=9,
            color=ACCENT5,
        )
        add_text_box(
            slide,
            term_x + Inches(3.6),
            try4 + Inches(0.05),
            Inches(0.6),
            Inches(0.24),
            c_count,
            font_size=9,
            color=TEXT,
            alignment=PP_ALIGN.CENTER,
        )
        add_text_box(
            slide,
            term_x + Inches(4.4),
            try4 + Inches(0.05),
            Inches(0.6),
            Inches(0.24),
            v_count,
            font_size=9,
            color=ACCENT5,
            alignment=PP_ALIGN.CENTER,
        )
        add_text_box(
            slide,
            term_x + Inches(5.1),
            try4 + Inches(0.05),
            Inches(0.8),
            Inches(0.24),
            action,
            font_size=9,
            color=MUTED,
            italic=True,
        )
        try4 += Inches(0.35)

    # --- Definition Consistency ---
    dfn_x, dfn_y = Inches(0.4), Inches(4.4)
    dfn_w, dfn_h = Inches(6.2), Inches(2.6)
    _, dfn_cy = mockup_frame(
        slide, dfn_x, dfn_y, dfn_w, dfn_h, "Definition Consistency  \u2014  --check-dfn"
    )

    add_text_box(
        slide,
        dfn_x + Inches(0.15),
        dfn_cy + Inches(0.1),
        Inches(3),
        Inches(0.22),
        "Undefined References (2)",
        font_size=10,
        bold=True,
        color=ACCENT5,
    )
    dfn_items_1 = [
        "\u2022  superblock referenced 4 times but never defined with <dfn>",
        "\u2022  palette_mode referenced 2 times but never defined with <dfn>",
    ]
    dy2 = dfn_cy + Inches(0.35)
    for item in dfn_items_1:
        add_text_box(
            slide,
            dfn_x + Inches(0.25),
            dy2,
            Inches(5.7),
            Inches(0.22),
            item,
            font_size=9,
            color=TEXT_LIGHT,
        )
        dy2 += Inches(0.22)

    dy2 += Inches(0.1)
    add_text_box(
        slide,
        dfn_x + Inches(0.15),
        dy2,
        Inches(3),
        Inches(0.22),
        "Unused Definitions (3)",
        font_size=10,
        bold=True,
        color=ACCENT3,
    )
    dy2 += Inches(0.25)
    dfn_items_2 = [
        "\u2022  <dfn>temporal unit</dfn> defined in \u00a72.1 but never referenced",
        "\u2022  <dfn>operating point</dfn> defined in \u00a73.2 but never referenced",
        "\u2022  <dfn>film grain unit</dfn> defined in \u00a77.4 but never referenced",
    ]
    for item in dfn_items_2:
        add_text_box(
            slide,
            dfn_x + Inches(0.25),
            dy2,
            Inches(5.7),
            Inches(0.22),
            item,
            font_size=9,
            color=TEXT_LIGHT,
        )
        dy2 += Inches(0.22)

    # --- Accessibility Audit ---
    a11y_x, a11y_y = Inches(6.8), Inches(4.4)
    a11y_w, a11y_h = Inches(6.1), Inches(2.6)
    _, a11y_cy = mockup_frame(
        slide, a11y_x, a11y_y, a11y_w, a11y_h, "Accessibility Audit  \u2014  --accessibility-audit"
    )

    a11y_items = [
        ("WCAG 1.1.1", "error", "Image missing alt text", "\u00a76.5 fig. 8"),
        ("WCAG 1.3.1", "warning", "Table header cell missing scope", "\u00a75.1"),
        ("WCAG 2.4.6", "warning", "Empty heading element", "\u00a78.2"),
        ("WCAG 4.1.2", "error", "Form control missing label", "\u00a7A.1"),
    ]

    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        a11y_x + Inches(0.1),
        a11y_cy + Inches(0.1),
        Inches(5.9),
        Inches(0.28),
        fill_color=NEAR_BLACK,
    )
    for cx2, ct in [
        (a11y_x + Inches(0.15), "Rule"),
        (a11y_x + Inches(1.2), "Level"),
        (a11y_x + Inches(2.0), "Issue"),
        (a11y_x + Inches(4.3), "Location"),
    ]:
        add_text_box(
            slide,
            cx2,
            a11y_cy + Inches(0.12),
            Inches(2.0),
            Inches(0.24),
            ct,
            font_size=9,
            bold=True,
            color=WHITE,
        )

    ary = a11y_cy + Inches(0.4)
    for i, (rule, level, issue, loc) in enumerate(a11y_items):
        bg = CARD_BG if i % 2 == 0 else WHITE
        add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            a11y_x + Inches(0.1),
            ary,
            Inches(5.9),
            Inches(0.35),
            fill_color=bg,
        )
        add_text_box(
            slide,
            a11y_x + Inches(0.15),
            ary + Inches(0.05),
            Inches(1.0),
            Inches(0.24),
            rule,
            font_size=9,
            bold=True,
            color=ACCENT,
            font_name="Consolas",
        )
        lv_color = ACCENT5 if level == "error" else ACCENT3
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            a11y_x + Inches(1.2),
            ary + Inches(0.05),
            Inches(0.65),
            Inches(0.24),
            fill_color=lv_color,
        )
        add_text_box(
            slide,
            a11y_x + Inches(1.2),
            ary + Inches(0.05),
            Inches(0.65),
            Inches(0.24),
            level,
            font_size=9,
            bold=True,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
        )
        add_text_box(
            slide,
            a11y_x + Inches(2.0),
            ary + Inches(0.05),
            Inches(2.2),
            Inches(0.24),
            issue,
            font_size=9,
            color=TEXT,
        )
        add_text_box(
            slide,
            a11y_x + Inches(4.3),
            ary + Inches(0.05),
            Inches(1.5),
            Inches(0.24),
            loc,
            font_size=9,
            color=TEXT_LIGHT,
        )
        ary += Inches(0.35)


def slide_05_enhancements(prs, layout):
    """Document Enhancements (18)."""
    slide = prs.slides.add_slide(layout)
    slide_title_bar(
        slide,
        "Document Enhancements (18)",
        "Sequential DOM mutations, ordered by priority  \u2022  Precompute where possible",
    )

    enhancements = [
        ("Equation Numbering", "--number-equations", "Auto-number (section.N)", ACCENT),
        ("Change Bars", "--change-bars [REF]", "Mark modified text", ACCENT),
        ("Revision History", "--revision-history", "Git tags \u2192 table", ACCENT),
        ("Table of Changes", "--table-of-changes [REF]", "List modified sections", ACCENT),
        ("Index Management", "--index [style]", "Alphabetical / remove", ACCENT),
        ("RFC 2119 Keywords", "--highlight-keywords", "Visual SHALL/MUST", ACCENT2),
        ("Fig/Table Tooltips", "--figure-table-tooltips", "Hover preview on xrefs", ACCENT2),
        ("Syntax Tooltips", "--syntax-tooltips", "SDL element info", ACCENT2),
        ("TOC Styling", "--toc-bold-primary-only", "Lightweight TOC", ACCENT2),
        ("Content Width", "--content-width [W]", "Constrain body width", ACCENT2),
        ("Line Anchors", "--line-anchors", "Deep-link code lines", ACCENT3),
        ("PWA", "--pwa", "Offline viewing", ACCENT3),
        ("Stability Badges", "--stability", "New/active markers", ACCENT3),
        ("Watermark", "--watermark [TEXT]", "Draft / confidential", ACCENT3),
        ("Cover Page", "--cover-page", "Styled title page", ACCENT4),
        ("Page Numbers", "--page-numbers [STYLE]", "Dual / arabic / none", ACCENT4),
    ]

    cols = 4
    card_w = Inches(2.95)
    card_h = Inches(1.0)
    start_x = Inches(0.45)
    start_y = Inches(1.5)

    for i, (title, flag, desc, color) in enumerate(enhancements):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + Inches(0.15))
        y = start_y + row * (card_h + Inches(0.1))

        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            y,
            card_w,
            card_h,
            fill_color=CARD_BG,
            line_color=BORDER,
            line_width=Pt(1),
        )
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, Pt(4), card_h, fill_color=color)
        add_text_box(
            slide,
            x + Inches(0.15),
            y + Inches(0.05),
            card_w - Inches(0.2),
            Inches(0.25),
            title,
            font_size=12,
            bold=True,
            color=color,
        )
        add_text_box(
            slide,
            x + Inches(0.15),
            y + Inches(0.3),
            card_w - Inches(0.2),
            Inches(0.22),
            flag,
            font_size=9,
            color=MUTED,
            font_name="Consolas",
        )
        add_text_box(
            slide,
            x + Inches(0.15),
            y + Inches(0.58),
            card_w - Inches(0.2),
            Inches(0.3),
            desc,
            font_size=10,
            color=TEXT_LIGHT,
        )

    add_text_box(
        slide,
        Inches(0.5),
        Inches(6.05),
        Inches(11),
        Inches(0.4),
        "Defaults enabled: --number-equations, --figure-table-tooltips, "
        "--syntax-tooltips, --toc-bold-primary-only",
        font_size=12,
        color=MUTED,
    )


def slide_06_tooltips(prs, layout):
    """Interactive Features: Tooltips & Search."""
    slide = prs.slides.add_slide(layout)
    slide_title_bar(
        slide,
        "Interactive Features: Tooltips & Search",
        "Hover previews and client-side search for improved navigation",
    )

    # --- Figure/Table Tooltip Mockup ---
    fx, fy = Inches(0.4), Inches(1.5)
    fw, fh = Inches(6.0), Inches(2.8)
    _, fc_y = mockup_frame(
        slide, fx, fy, fw, fh, "Figure/Table Tooltips  \u2014  --figure-table-tooltips"
    )

    add_text_box(
        slide,
        fx + Inches(0.2),
        fc_y + Inches(0.15),
        Inches(5.5),
        Inches(0.3),
        "6.3.2 Loop Filter Process",
        font_size=12,
        bold=True,
        color=NEAR_BLACK,
    )

    prose_y = fc_y + Inches(0.5)
    add_rich_text(
        slide,
        fx + Inches(0.2),
        prose_y,
        Inches(5.5),
        Inches(0.5),
        [
            ("The loop filter process applies the filter shown in ", TEXT),
            ("Table 12", ACCENT, True),
            (" to each", TEXT),
        ],
        font_size=10,
    )
    add_rich_text(
        slide,
        fx + Inches(0.2),
        prose_y + Inches(0.3),
        Inches(5.5),
        Inches(0.5),
        [
            ("reconstructed sample. The filter taps are illustrated in ", TEXT),
            ("Figure 7", ACCENT, True),
            (".", TEXT),
        ],
        font_size=10,
    )

    # Tooltip popup
    tt_x = fx + Inches(2.0)
    tt_y = prose_y + Inches(0.7)
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        tt_x,
        tt_y,
        Inches(3.5),
        Inches(1.2),
        fill_color=NEAR_BLACK,
        line_color=ACCENT,
        line_width=Pt(2),
    )
    add_shape(
        slide,
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        tt_x + Inches(0.8),
        tt_y - Inches(0.15),
        Inches(0.2),
        Inches(0.15),
        fill_color=NEAR_BLACK,
    )
    add_text_box(
        slide,
        tt_x + Inches(0.12),
        tt_y + Inches(0.08),
        Inches(3.3),
        Inches(0.22),
        "Table 12 \u2014 Loop Filter Parameters",
        font_size=9,
        bold=True,
        color=ACCENT,
    )
    add_text_box(
        slide,
        tt_x + Inches(0.12),
        tt_y + Inches(0.32),
        Inches(3.3),
        Inches(0.7),
        "Defines the filter strength and taps\nfor each block size "
        "configuration.\nSection 6.3 \u2022 Page 42",
        font_size=9,
        color=RGBColor(0xCB, 0xD5, 0xE0),
    )

    add_shape(
        slide,
        MSO_SHAPE.OVAL,
        fx + Inches(3.55),
        prose_y + Inches(0.32),
        Inches(0.08),
        Inches(0.08),
        fill_color=ACCENT,
    )
    add_text_box(
        slide,
        fx + Inches(3.7),
        prose_y + Inches(0.25),
        Inches(1.0),
        Inches(0.2),
        "\u2190 hover",
        font_size=9,
        color=ACCENT,
        italic=True,
    )

    # --- Syntax Tooltip Mockup ---
    sx, sy_start = Inches(6.8), Inches(1.5)
    sw, sh = Inches(6.1), Inches(2.8)
    _, sc_y = mockup_frame(
        slide, sx, sy_start, sw, sh, "Syntax Tooltips  \u2014  --syntax-tooltips"
    )

    add_text_box(
        slide,
        sx + Inches(0.15),
        sc_y + Inches(0.1),
        Inches(5.5),
        Inches(0.25),
        "frame_header( ) {",
        font_size=11,
        bold=True,
        color=NEAR_BLACK,
        font_name="Consolas",
    )
    sdl_rows = [
        ("  show_existing_frame", "f(1)"),
        ("  frame_type", "f(2)"),
        ("  show_frame", "f(1)"),
        ("  error_resilient_mode", "f(1)"),
    ]
    ry = sc_y + Inches(0.4)
    for name, desc in sdl_rows:
        add_text_box(
            slide,
            sx + Inches(0.15),
            ry,
            Inches(2.8),
            Inches(0.22),
            name,
            font_size=9,
            color=ACCENT,
            font_name="Consolas",
        )
        add_text_box(
            slide,
            sx + Inches(3.2),
            ry,
            Inches(1.5),
            Inches(0.22),
            desc,
            font_size=9,
            color=MUTED,
            font_name="Consolas",
        )
        ry += Inches(0.24)

    stt_x = sx + Inches(1.5)
    stt_y = sc_y + Inches(0.75)
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        stt_x,
        stt_y,
        Inches(3.8),
        Inches(1.0),
        fill_color=NEAR_BLACK,
        line_color=ACCENT4,
        line_width=Pt(2),
    )
    add_text_box(
        slide,
        stt_x + Inches(0.1),
        stt_y + Inches(0.06),
        Inches(3.5),
        Inches(0.2),
        "frame_type",
        font_size=9,
        bold=True,
        color=ACCENT4,
        font_name="Consolas",
    )
    add_text_box(
        slide,
        stt_x + Inches(0.1),
        stt_y + Inches(0.28),
        Inches(3.5),
        Inches(0.6),
        "Specifies the type of this frame.\n0 = KEY_FRAME, 1 = INTER_FRAME,\n"
        "2 = INTRA_ONLY, 3 = S_FRAME",
        font_size=9,
        color=RGBColor(0xCB, 0xD5, 0xE0),
    )

    # --- Search Overlay Mockup ---
    search_x, search_y = Inches(0.4), Inches(4.7)
    search_w, search_h = Inches(6.0), Inches(2.4)
    _, search_cy = mockup_frame(
        slide, search_x, search_y, search_w, search_h, "Client-Side Search  \u2014  --search-index"
    )

    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        search_x + Inches(0.3),
        search_cy + Inches(0.15),
        Inches(5.3),
        Inches(0.4),
        fill_color=WHITE,
        line_color=ACCENT,
        line_width=Pt(2),
    )
    add_text_box(
        slide,
        search_x + Inches(0.45),
        search_cy + Inches(0.2),
        Inches(4.5),
        Inches(0.3),
        "\U0001f50d  loop filter",
        font_size=11,
        color=TEXT,
    )
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        search_x + Inches(4.7),
        search_cy + Inches(0.2),
        Inches(0.7),
        Inches(0.25),
        fill_color=LIGHT_BG,
        line_color=BORDER,
        line_width=Pt(1),
    )
    add_text_box(
        slide,
        search_x + Inches(4.7),
        search_cy + Inches(0.2),
        Inches(0.7),
        Inches(0.25),
        "\u2318K",
        font_size=9,
        bold=True,
        color=MUTED,
        alignment=PP_ALIGN.CENTER,
    )

    results = [
        ("\u00a76.3", "Loop Filter Process", "The loop filter process shall be invoked..."),
        ("\u00a76.3.2", "Loop Filter Parameters", "Table 12 defines the filter strength..."),
        ("\u00a76.3.5", "Deblocking Filter", "The deblocking loop filter operates on..."),
    ]
    ry = search_cy + Inches(0.65)
    for sec, title, snippet in results:
        add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            search_x + Inches(0.3),
            ry,
            Inches(5.3),
            Inches(0.45),
            fill_color=CARD_BG,
            line_color=BORDER,
            line_width=Pt(1),
        )
        add_text_box(
            slide,
            search_x + Inches(0.4),
            ry + Inches(0.02),
            Inches(0.6),
            Inches(0.2),
            sec,
            font_size=9,
            bold=True,
            color=ACCENT,
        )
        add_text_box(
            slide,
            search_x + Inches(1.0),
            ry + Inches(0.02),
            Inches(4.0),
            Inches(0.2),
            title,
            font_size=9,
            bold=True,
            color=TEXT,
        )
        add_text_box(
            slide,
            search_x + Inches(0.4),
            ry + Inches(0.22),
            Inches(5.0),
            Inches(0.2),
            snippet,
            font_size=9,
            color=TEXT_LIGHT,
        )
        ry += Inches(0.5)

    # --- Keyword Highlighting Mockup ---
    kw_x, kw_y = Inches(6.8), Inches(4.7)
    kw_w, kw_h = Inches(6.1), Inches(2.4)
    _, kw_cy = mockup_frame(
        slide, kw_x, kw_y, kw_w, kw_h, "RFC 2119 Keywords  \u2014  --highlight-keywords"
    )

    add_text_box(
        slide,
        kw_x + Inches(0.15),
        kw_cy + Inches(0.1),
        Inches(5.5),
        Inches(0.25),
        "6.2.1 General Requirements",
        font_size=11,
        bold=True,
        color=NEAR_BLACK,
    )

    # Draw all keyword highlight backgrounds FIRST (so text renders on top)
    ky = kw_cy + Inches(0.42)
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        kw_x + Inches(0.88),
        ky - Inches(0.02),
        Inches(0.55),
        Inches(0.24),
        fill_color=RGBColor(0xFE, 0xE2, 0xE2),
    )
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        kw_x + Inches(2.3),
        ky + Inches(0.33),
        Inches(0.95),
        Inches(0.24),
        fill_color=RGBColor(0xFE, 0xE2, 0xE2),
    )
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        kw_x + Inches(0.88),
        ky + Inches(0.68),
        Inches(0.82),
        Inches(0.24),
        fill_color=RGBColor(0xFF, 0xED, 0xD5),
    )
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        kw_x + Inches(1.4),
        ky + Inches(1.03),
        Inches(0.48),
        Inches(0.24),
        fill_color=RGBColor(0xDB, 0xED, 0xF9),
    )

    # Now draw the text on top
    add_rich_text(
        slide,
        kw_x + Inches(0.15),
        ky,
        Inches(5.7),
        Inches(0.3),
        [
            ("Decoders ", TEXT),
            ("SHALL", KW_SHALL, True),
            (" support all profiles defined in Table 3.", TEXT),
        ],
        font_size=10,
    )

    ky += Inches(0.35)
    add_rich_text(
        slide,
        kw_x + Inches(0.15),
        ky,
        Inches(5.7),
        Inches(0.3),
        [("The value of frame_type ", TEXT), ("MUST NOT", KW_SHALL, True), (" exceed 3.", TEXT)],
        font_size=10,
    )

    ky += Inches(0.35)
    add_rich_text(
        slide,
        kw_x + Inches(0.15),
        ky,
        Inches(5.7),
        Inches(0.3),
        [
            ("Encoders ", TEXT),
            ("SHOULD", KW_SHOULD, True),
            (" signal the error resilient mode flag.", TEXT),
        ],
        font_size=10,
    )

    ky += Inches(0.35)
    add_rich_text(
        slide,
        kw_x + Inches(0.15),
        ky,
        Inches(5.7),
        Inches(0.3),
        [
            ("Implementations ", TEXT),
            ("MAY", KW_MAY, True),
            (" use alternative decoding orders.", TEXT),
        ],
        font_size=10,
    )

    ky += Inches(0.45)
    legend_items = [
        ("SHALL / MUST", KW_SHALL, RGBColor(0xFE, 0xE2, 0xE2)),
        ("SHOULD", KW_SHOULD, RGBColor(0xFF, 0xED, 0xD5)),
        ("MAY", KW_MAY, RGBColor(0xDB, 0xED, 0xF9)),
    ]
    lx = kw_x + Inches(0.3)
    for word, color, bg in legend_items:
        add_shape(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE, lx, ky, Inches(0.12), Inches(0.12), fill_color=bg
        )
        add_text_box(
            slide,
            lx + Inches(0.18),
            ky - Inches(0.02),
            Inches(1.2),
            Inches(0.2),
            word,
            font_size=9,
            bold=True,
            color=color,
            font_name="Consolas",
        )
        lx += Inches(1.6)


def slide_07_change_tracking(prs, layout):
    """Change Tracking: Change Bars & Diff Viewer."""
    slide = prs.slides.add_slide(layout)
    slide_title_bar(
        slide,
        "Change Tracking: Change Bars & Diff Viewer",
        "Visual indicators for modified content since a baseline reference",
    )

    # --- Change Bars Mockup ---
    cb_x, cb_y = Inches(0.4), Inches(1.5)
    cb_w, cb_h = Inches(6.0), Inches(5.3)
    _, cb_cy = mockup_frame(
        slide, cb_x, cb_y, cb_w, cb_h, "Change Bars  \u2014  --change-bars [REF]"
    )

    add_text_box(
        slide,
        cb_x + Inches(0.35),
        cb_cy + Inches(0.1),
        Inches(5.3),
        Inches(0.25),
        "6.2 Frame Header Semantics",
        font_size=12,
        bold=True,
        color=NEAR_BLACK,
    )

    cb_ty = cb_cy + Inches(0.45)
    add_text_box(
        slide,
        cb_x + Inches(0.35),
        cb_ty,
        Inches(5.3),
        Inches(0.3),
        "The frame header contains information about the coding",
        font_size=10,
        color=TEXT,
    )
    cb_ty += Inches(0.28)
    add_text_box(
        slide,
        cb_x + Inches(0.35),
        cb_ty,
        Inches(5.3),
        Inches(0.3),
        "parameters used for the current frame.",
        font_size=10,
        color=TEXT,
    )

    cb_ty += Inches(0.4)
    for line_text in [
        "When error_resilient_mode is equal to 1, the decoder",
        "shall reset all reference frame state at the start of",
        "each tile group. This ensures independent decoding",
        "of tile groups for improved error recovery.",
    ]:
        add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            cb_x + Inches(0.15),
            cb_ty,
            Pt(4),
            Inches(0.26),
            fill_color=ACCENT2,
        )
        add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            cb_x + Inches(0.25),
            cb_ty,
            Inches(5.4),
            Inches(0.26),
            fill_color=CB_GREEN,
        )
        add_text_box(
            slide,
            cb_x + Inches(0.35),
            cb_ty,
            Inches(5.3),
            Inches(0.26),
            line_text,
            font_size=10,
            color=CB_GREEN_DARK,
        )
        cb_ty += Inches(0.26)

    cb_ty += Inches(0.15)
    add_text_box(
        slide,
        cb_x + Inches(0.35),
        cb_ty,
        Inches(5.3),
        Inches(0.3),
        "The frame header is followed by the tile group header",
        font_size=10,
        color=TEXT,
    )
    cb_ty += Inches(0.28)
    add_text_box(
        slide,
        cb_x + Inches(0.35),
        cb_ty,
        Inches(5.3),
        Inches(0.3),
        "which specifies the tile decoding order.",
        font_size=10,
        color=TEXT,
    )

    cb_ty += Inches(0.4)
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        cb_x + Inches(0.15),
        cb_ty,
        Pt(4),
        Inches(0.26),
        fill_color=ACCENT3,
    )
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        cb_x + Inches(0.25),
        cb_ty,
        Inches(5.4),
        Inches(0.26),
        fill_color=CB_YELLOW,
    )
    add_text_box(
        slide,
        cb_x + Inches(0.35),
        cb_ty,
        Inches(5.3),
        Inches(0.26),
        "The maximum number of tile groups is 256 (changed from 128).",
        font_size=10,
        color=CB_YELLOW_DARK,
    )

    cb_ty += Inches(0.55)
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        cb_x + Inches(0.3),
        cb_ty,
        Pt(4),
        Inches(0.18),
        fill_color=ACCENT2,
    )
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        cb_x + Inches(0.4),
        cb_ty,
        Inches(0.5),
        Inches(0.18),
        fill_color=CB_GREEN,
    )
    add_text_box(
        slide,
        cb_x + Inches(1.0),
        cb_ty - Inches(0.02),
        Inches(0.8),
        Inches(0.2),
        "Added",
        font_size=9,
        color=ACCENT2,
    )
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        cb_x + Inches(2.0),
        cb_ty,
        Pt(4),
        Inches(0.18),
        fill_color=ACCENT3,
    )
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        cb_x + Inches(2.1),
        cb_ty,
        Inches(0.5),
        Inches(0.18),
        fill_color=CB_YELLOW,
    )
    add_text_box(
        slide,
        cb_x + Inches(2.7),
        cb_ty - Inches(0.02),
        Inches(0.9),
        Inches(0.2),
        "Modified",
        font_size=9,
        color=ACCENT3,
    )

    # --- Table of Changes Mockup ---
    tc_x, tc_y = Inches(6.8), Inches(1.5)
    tc_w, tc_h = Inches(6.1), Inches(2.5)
    _, tcc_y = mockup_frame(
        slide, tc_x, tc_y, tc_w, tc_h, "Table of Changes  \u2014  --table-of-changes [REF]"
    )

    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        tc_x + Inches(0.1),
        tcc_y + Inches(0.1),
        Inches(5.9),
        Inches(0.35),
        fill_color=NEAR_BLACK,
    )
    for col_x, col_text, col_w in [
        (tc_x + Inches(0.2), "Section", Inches(1.0)),
        (tc_x + Inches(1.3), "Title", Inches(2.5)),
        (tc_x + Inches(3.9), "Lines Changed", Inches(0.9)),
        (tc_x + Inches(4.9), "Status", Inches(0.9)),
    ]:
        add_text_box(
            slide,
            col_x,
            tcc_y + Inches(0.12),
            col_w,
            Inches(0.3),
            col_text,
            font_size=9,
            bold=True,
            color=WHITE,
        )

    toc_rows = [
        ("6.2", "Frame Header", "+42 / -8", "Modified", ACCENT3),
        ("6.3", "Loop Filter Process", "+18 / -3", "Modified", ACCENT3),
        ("6.5", "CDEF Process", "+95 / -0", "New", ACCENT2),
        ("A.3", "Decoder Model", "+12 / -12", "Modified", ACCENT3),
    ]
    try_ = tcc_y + Inches(0.5)
    for i, (sec, title, lines, status, color) in enumerate(toc_rows):
        bg = CARD_BG if i % 2 == 0 else WHITE
        add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            tc_x + Inches(0.1),
            try_,
            Inches(5.9),
            Inches(0.35),
            fill_color=bg,
        )
        add_text_box(
            slide,
            tc_x + Inches(0.2),
            try_ + Inches(0.05),
            Inches(0.9),
            Inches(0.25),
            sec,
            font_size=9,
            bold=True,
            color=ACCENT,
        )
        add_text_box(
            slide,
            tc_x + Inches(1.3),
            try_ + Inches(0.05),
            Inches(2.3),
            Inches(0.25),
            title,
            font_size=9,
            color=TEXT,
        )
        add_text_box(
            slide,
            tc_x + Inches(3.9),
            try_ + Inches(0.05),
            Inches(0.9),
            Inches(0.25),
            lines,
            font_size=9,
            color=TEXT_LIGHT,
            font_name="Consolas",
        )
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            tc_x + Inches(4.9),
            try_ + Inches(0.05),
            Inches(0.8),
            Inches(0.24),
            fill_color=color,
        )
        add_text_box(
            slide,
            tc_x + Inches(4.9),
            try_ + Inches(0.05),
            Inches(0.8),
            Inches(0.24),
            status,
            font_size=9,
            bold=True,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
        )
        try_ += Inches(0.35)

    # --- Revision History Mockup ---
    rh_x, rh_y = Inches(6.8), Inches(4.3)
    rh_w, rh_h = Inches(6.1), Inches(2.5)
    _, rh_cy = mockup_frame(
        slide, rh_x, rh_y, rh_w, rh_h, "Revision History  \u2014  --revision-history"
    )

    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        rh_x + Inches(0.1),
        rh_cy + Inches(0.1),
        Inches(5.9),
        Inches(0.32),
        fill_color=NEAR_BLACK,
    )
    for col_x, col_text in [
        (rh_x + Inches(0.2), "Version"),
        (rh_x + Inches(1.2), "Date"),
        (rh_x + Inches(2.3), "Author"),
        (rh_x + Inches(3.5), "Description"),
    ]:
        add_text_box(
            slide,
            col_x,
            rh_cy + Inches(0.12),
            Inches(1.2),
            Inches(0.28),
            col_text,
            font_size=9,
            bold=True,
            color=WHITE,
        )

    rev_rows = [
        ("v0.4", "2026-03-15", "A. Tourapis", "Phase 6: DecCT formula fix"),
        ("v0.3", "2026-03-10", "A. Tourapis", "Phase 5: VBI fix and clarifications"),
        ("v0.2", "2026-03-05", "A. Tourapis", "Phase 4: Syntax element corrections"),
        ("v0.1", "2026-02-28", "A. Tourapis", "Initial working draft"),
    ]
    rry = rh_cy + Inches(0.45)
    for i, (ver, date, author, desc) in enumerate(rev_rows):
        bg = CARD_BG if i % 2 == 0 else WHITE
        add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            rh_x + Inches(0.1),
            rry,
            Inches(5.9),
            Inches(0.32),
            fill_color=bg,
        )
        add_text_box(
            slide,
            rh_x + Inches(0.2),
            rry + Inches(0.04),
            Inches(0.9),
            Inches(0.24),
            ver,
            font_size=9,
            bold=True,
            color=ACCENT4,
        )
        add_text_box(
            slide,
            rh_x + Inches(1.2),
            rry + Inches(0.04),
            Inches(1.0),
            Inches(0.24),
            date,
            font_size=9,
            color=TEXT_LIGHT,
        )
        add_text_box(
            slide,
            rh_x + Inches(2.3),
            rry + Inches(0.04),
            Inches(1.1),
            Inches(0.24),
            author,
            font_size=9,
            color=TEXT,
        )
        add_text_box(
            slide,
            rh_x + Inches(3.5),
            rry + Inches(0.04),
            Inches(2.3),
            Inches(0.24),
            desc,
            font_size=9,
            color=TEXT_LIGHT,
        )
        rry += Inches(0.32)


def slide_08_output_tasks(prs, layout):
    """Output Tasks (20)."""
    slide = prs.slides.add_slide(layout)
    slide_title_bar(
        slide,
        "Output Tasks (20)",
        "Independent tasks \u2014 run concurrently with --parallel-outputs",
    )

    output_groups = [
        (
            "Document Formats",
            ACCENT,
            [
                ("PDF", "--pdf / --weasyprint"),
                ("PDF/A", "--pdfa"),
                ("Word", "--docx"),
                ("LaTeX", "--latex"),
                ("Standalone HTML", "--standalone"),
            ],
        ),
        (
            "Analysis & Reports",
            ACCENT2,
            [
                ("Spec Metrics", "--spec-metrics"),
                ("Cross-Ref Report", "--xref-report"),
                ("Compliance Matrix", "--compliance-matrix"),
                ("Normative Deps", "--normative-deps"),
                ("Attribution", "--attribution"),
            ],
        ),
        (
            "Change Tracking",
            ACCENT3,
            [
                ("Change Summary", "--change-summary"),
                ("PR Summary", "--pr-summary"),
                ("Regression Check", "--regression"),
                ("Save Baseline", "--save-baseline"),
                ("Spec Compare", "--spec-compare"),
            ],
        ),
        (
            "Standards Workflow",
            ACCENT4,
            [
                ("Requirement IDs", "--requirement-ids"),
                ("Issue Traceability", "--issue-traceability"),
                ("Ext Spec Deps", "--ext-spec-deps"),
                ("Metrics Trend", "--metrics-trend"),
                ("Release Automation", "--release"),
            ],
        ),
    ]

    group_w = Inches(3.0)
    start_x = Inches(0.3)
    start_y = Inches(1.5)

    for gi, (group_name, color, items) in enumerate(output_groups):
        gx = start_x + gi * (group_w + Inches(0.15))
        add_shape(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE, gx, start_y, group_w, Inches(0.42), fill_color=color
        )
        add_text_box(
            slide,
            gx,
            start_y + Inches(0.06),
            group_w,
            Inches(0.3),
            group_name,
            font_size=13,
            bold=True,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
        )
        for ii, (name, flag) in enumerate(items):
            iy = start_y + Inches(0.52) + ii * Inches(0.85)
            add_shape(
                slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                gx,
                iy,
                group_w,
                Inches(0.75),
                fill_color=CARD_BG,
                line_color=BORDER,
                line_width=Pt(1),
            )
            add_text_box(
                slide,
                gx + Inches(0.12),
                iy + Inches(0.08),
                group_w - Inches(0.24),
                Inches(0.25),
                name,
                font_size=12,
                bold=True,
                color=NEAR_BLACK,
            )
            add_text_box(
                slide,
                gx + Inches(0.12),
                iy + Inches(0.38),
                group_w - Inches(0.24),
                Inches(0.25),
                flag,
                font_size=9,
                color=MUTED,
                font_name="Consolas",
            )


def slide_09_profiles(prs, layout):
    """Build Profiles & Configuration."""
    slide = prs.slides.add_slide(layout)
    slide_title_bar(
        slide,
        "Build Profiles & Configuration",
        "Predefined flag sets for common workflows  \u2022  Extensible via TOML",
    )

    profiles = [
        ("quick", "Fast HTML-only build\nNo enhancements", "(default flags)", ACCENT),
        (
            "draft",
            "Working draft\nEquations + keywords + validation",
            "number_equations, highlight_keywords,\nvalidate_refs, tooltips",
            ACCENT2,
        ),
        (
            "review",
            "Review build\nAll quality checks + change bars",
            "draft + check_terminology,\ncheck_orphan_refs, change_bars",
            ACCENT3,
        ),
        (
            "publication",
            "Full publication build\nPDF, LOF, LOT, all enhancements",
            "review + revision_history,\npdf, lof, lot, pwa",
            ACCENT4,
        ),
        (
            "pdf-draft",
            "Quick PDF generation\nMinimal enhancements",
            "number_equations, pdf",
            ACCENT5,
        ),
        (
            "pdf-final",
            "Final PDF\nAll bells and whistles",
            "pdf-draft + highlight_keywords,\nrevision_history, lof, lot",
            TEAL,
        ),
    ]

    card_w = Inches(3.9)
    card_h = Inches(2.2)

    for i, (name, desc, flags, color) in enumerate(profiles):
        col = i % 3
        row = i // 3
        x = Inches(0.4) + col * (card_w + Inches(0.2))
        y = Inches(1.55) + row * (card_h + Inches(0.2))

        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            y,
            card_w,
            card_h,
            fill_color=CARD_BG,
            line_color=color,
            line_width=Pt(2),
        )

        badge_w = Inches(2.1)
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches(0.15),
            y + Inches(0.12),
            badge_w,
            Inches(0.35),
            fill_color=color,
        )
        add_text_box(
            slide,
            x + Inches(0.15),
            y + Inches(0.12),
            badge_w,
            Inches(0.35),
            f"--profile {name}",
            font_size=11,
            bold=True,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
            font_name="Consolas",
        )

        add_text_box(
            slide,
            x + Inches(0.15),
            y + Inches(0.58),
            card_w - Inches(0.3),
            Inches(0.65),
            desc,
            font_size=12,
            color=TEXT,
        )
        add_text_box(
            slide,
            x + Inches(0.15),
            y + Inches(1.35),
            card_w - Inches(0.3),
            Inches(0.75),
            flags,
            font_size=9,
            color=MUTED,
            font_name="Consolas",
        )

    add_text_box(
        slide,
        Inches(0.4),
        Inches(6.25),
        Inches(12),
        Inches(0.4),
        "Custom profiles: --profiles-file profiles.toml  \u2022  "
        "Project config: specbuild.toml or [tool.specbuild] in pyproject.toml",
        font_size=12,
        color=MUTED,
    )


def slide_10_pdf(prs, layout):
    """PDF Generation."""
    slide = prs.slides.add_slide(layout)
    slide_title_bar(
        slide,
        "PDF Generation",
        "Two rendering engines with full control over layout and typography",
    )

    engines = [
        (
            "Chrome / Headless",
            ACCENT,
            [
                "Default engine (--pdf)",
                "Best CSS compatibility for complex layouts",
                "Requires Chrome or Chromium installed",
                "Ideal for pixel-perfect output",
            ],
        ),
        (
            "WeasyPrint",
            ACCENT2,
            [
                "Pure Python engine (--weasyprint)",
                "No browser dependency required",
                "pip install specbuild[pdf]",
                "Best for CI pipelines and automation",
            ],
        ),
    ]

    for i, (title, color, items) in enumerate(engines):
        x = Inches(0.5) + i * Inches(4.2)
        add_card(
            slide,
            x,
            Inches(1.6),
            Inches(3.9),
            Inches(2.2),
            title,
            items,
            accent_color=color,
            title_size=16,
            body_size=12,
        )

    pdf_opts = [
        ("Page Size", "--page-size", "letter / a4 / legal"),
        ("Font Increase", "--pdf-font-increase N", "+1 to +3 points"),
        ("TOC Leaders", "--toc-leaders", "css / table / none"),
        ("List of Figures", "--lof", "With page numbers"),
        ("List of Tables", "--lot", "With page numbers"),
        ("Equation Font", "--equation-font", "6 font choices"),
        ("Equation Scale", "--equation-scale", "Relative to text"),
        ("Section Headers", "--section-headers", "Running titles"),
        ("Cover Page", "--cover-page", "Styled title page"),
        ("Page Numbers", "--page-numbers", "dual / arabic / none"),
    ]

    opt_x = Inches(8.8)
    add_text_box(
        slide,
        opt_x,
        Inches(1.6),
        Inches(4),
        Inches(0.35),
        "PDF OPTIONS",
        font_size=14,
        bold=True,
        color=ACCENT3,
    )

    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        opt_x,
        Inches(2.0),
        Inches(4.2),
        Inches(0.35),
        fill_color=NEAR_BLACK,
    )
    add_text_box(
        slide,
        opt_x + Inches(0.1),
        Inches(2.03),
        Inches(1.5),
        Inches(0.3),
        "Option",
        font_size=10,
        bold=True,
        color=WHITE,
    )
    add_text_box(
        slide,
        opt_x + Inches(2.2),
        Inches(2.03),
        Inches(1.8),
        Inches(0.3),
        "Values",
        font_size=10,
        bold=True,
        color=WHITE,
    )

    for i, (name, flag, desc) in enumerate(pdf_opts):
        iy = Inches(2.4) + i * Inches(0.38)
        bg = CARD_BG if i % 2 == 0 else WHITE
        add_shape(slide, MSO_SHAPE.RECTANGLE, opt_x, iy, Inches(4.2), Inches(0.38), fill_color=bg)
        add_text_box(
            slide,
            opt_x + Inches(0.1),
            iy + Inches(0.05),
            Inches(2.0),
            Inches(0.25),
            name,
            font_size=10,
            bold=True,
            color=TEXT,
        )
        add_text_box(
            slide,
            opt_x + Inches(2.2),
            iy + Inches(0.05),
            Inches(1.8),
            Inches(0.25),
            desc,
            font_size=10,
            color=TEXT_LIGHT,
        )


def slide_11_pdf_enhancements(prs, layout):
    """PDF Enhancements: Cover Page, TOC & Page Numbers."""
    slide = prs.slides.add_slide(layout)
    slide_title_bar(
        slide,
        "PDF Enhancements: Cover Page, TOC & Page Numbers",
        "Publication-ready PDF with front matter, leader dots, and dual numbering",
    )

    # --- Cover Page Mockup (wider to avoid overflow) ---
    cx, cy = Inches(0.4), Inches(1.5)
    cw, ch = Inches(3.8), Inches(5.3)
    _, content_y = mockup_frame(slide, cx, cy, cw, ch, "Cover Page  \u2014  --cover-page")

    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        cx + Inches(0.3),
        content_y + Inches(0.3),
        Inches(3.2),
        Pt(3),
        fill_color=ACCENT,
    )
    add_text_box(
        slide,
        cx + Inches(0.3),
        content_y + Inches(0.5),
        Inches(3.2),
        Inches(0.3),
        "Alliance for Open Media",
        font_size=10,
        color=MUTED,
        alignment=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        cx + Inches(0.3),
        content_y + Inches(1.0),
        Inches(3.2),
        Inches(0.8),
        "AV2 Bitstream &\nDecoding Process\nSpecification",
        font_size=16,
        bold=True,
        color=NEAR_BLACK,
        alignment=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        cx + Inches(0.3),
        content_y + Inches(2.0),
        Inches(3.2),
        Inches(0.3),
        "Working Draft",
        font_size=12,
        color=ACCENT,
        italic=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        cx + Inches(0.8),
        content_y + Inches(2.5),
        Inches(2.2),
        Pt(1),
        fill_color=BORDER,
    )
    add_text_box(
        slide,
        cx + Inches(0.3),
        content_y + Inches(2.7),
        Inches(3.2),
        Inches(0.6),
        "Document: AOMedia-AV2-001\nDate: March 15, 2026\nmain@fae537e",
        font_size=9,
        color=MUTED,
        alignment=PP_ALIGN.CENTER,
    )

    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        cx + Inches(2.8),
        content_y + Inches(4.2),
        Inches(0.6),
        Inches(0.25),
        fill_color=LIGHT_BG,
        line_color=ACCENT,
        line_width=Pt(1),
    )
    add_text_box(
        slide,
        cx + Inches(2.8),
        content_y + Inches(4.2),
        Inches(0.6),
        Inches(0.25),
        "i",
        font_size=9,
        color=ACCENT,
        alignment=PP_ALIGN.CENTER,
        font_name="Times New Roman",
        italic=True,
    )

    add_text_box(
        slide,
        cx,
        cy + ch + Inches(0.1),
        cw,
        Inches(0.3),
        "--cover-title, --cover-subtitle, --cover-organization",
        font_size=9,
        color=MUTED,
        alignment=PP_ALIGN.CENTER,
        font_name="Consolas",
    )

    # --- TOC Mockup ---
    tx, ty = Inches(4.5), Inches(1.5)
    tw, th = Inches(4.3), Inches(5.3)
    _, tc_y = mockup_frame(slide, tx, ty, tw, th, "Table of Contents  \u2014  --toc-leaders css")

    toc_entries = [
        ("1", "Scope", "1", True),
        ("2", "Normative References", "2", True),
        ("3", "Definitions", "3", True),
        ("4", "Conventions", "5", True),
        ("  4.1", "Arithmetic Operators", "5", False),
        ("  4.2", "Logical Operators", "6", False),
        ("  4.3", "Relational Operators", "7", False),
        ("5", "Syntax Structures", "8", True),
        ("  5.1", "Open Bitstream Unit", "8", False),
        ("  5.2", "Frame Header", "10", False),
        ("6", "Decoding Process", "15", True),
        ("  6.1", "General", "15", False),
        ("  6.2", "Tile Decoding", "18", False),
        ("  6.3", "Loop Filtering", "22", False),
    ]

    for i, (num, title, page, is_bold) in enumerate(toc_entries):
        ey = tc_y + Inches(0.15) + i * Inches(0.32)
        add_text_box(
            slide,
            tx + Inches(0.15),
            ey,
            Inches(0.5),
            Inches(0.25),
            num,
            font_size=9,
            bold=is_bold,
            color=TEXT if is_bold else TEXT_LIGHT,
        )
        add_text_box(
            slide,
            tx + Inches(0.6),
            ey,
            Inches(2.5),
            Inches(0.25),
            title,
            font_size=9,
            bold=is_bold,
            color=ACCENT if is_bold else TEXT,
        )
        add_text_box(
            slide,
            tx + Inches(2.8),
            ey,
            Inches(1.0),
            Inches(0.25),
            "\u00b7 " * 12,
            font_size=9,
            color=MUTED,
            alignment=PP_ALIGN.LEFT,
        )
        add_text_box(
            slide,
            tx + Inches(3.7),
            ey,
            Inches(0.4),
            Inches(0.25),
            page,
            font_size=9,
            bold=is_bold,
            color=TEXT,
            alignment=PP_ALIGN.RIGHT,
        )

    add_text_box(
        slide,
        tx,
        ty + th + Inches(0.1),
        tw,
        Inches(0.3),
        "--toc-leaders css|table|none  \u2022  --lof  \u2022  --lot",
        font_size=9,
        color=MUTED,
        alignment=PP_ALIGN.CENTER,
        font_name="Consolas",
    )

    # --- Page Number Styles ---
    px, py = Inches(9.1), Inches(1.5)
    pw, ph = Inches(3.8), Inches(5.3)
    _, pc_y = mockup_frame(slide, px, py, pw, ph, "Page Numbering  \u2014  --page-numbers")

    styles = [
        (
            "dual (default)",
            [
                ("Front matter:", "i, ii, iii, iv, v \u2026", "Roman"),
                ("Body:", "1, 2, 3, 4, 5 \u2026", "Arabic"),
                ("Reset:", "Page 1 at first section", ""),
            ],
        ),
        (
            "arabic",
            [
                ("Throughout:", "1, 2, 3, 4, 5 \u2026", "Sequential"),
            ],
        ),
        (
            "none",
            [
                ("No page numbers", "", "Clean output"),
            ],
        ),
    ]

    sy = pc_y + Inches(0.15)
    for style_name, items in styles:
        add_label(slide, px + Inches(0.2), sy, Inches(1.6), style_name, ACCENT4, font_size=9)
        sy += Inches(0.32)
        for label, value, note in items:
            add_text_box(
                slide,
                px + Inches(0.3),
                sy,
                Inches(1.1),
                Inches(0.2),
                label,
                font_size=9,
                bold=True,
                color=TEXT,
            )
            add_text_box(
                slide,
                px + Inches(1.4),
                sy,
                Inches(1.4),
                Inches(0.2),
                value,
                font_size=9,
                color=ACCENT,
                font_name="Consolas",
            )
            if note:
                add_text_box(
                    slide,
                    px + Inches(2.8),
                    sy,
                    Inches(0.8),
                    Inches(0.2),
                    note,
                    font_size=9,
                    color=MUTED,
                    italic=True,
                )
            sy += Inches(0.24)
        sy += Inches(0.1)

    # LOF/LOT
    sy += Inches(0.05)
    add_text_box(
        slide,
        px + Inches(0.2),
        sy,
        Inches(3.3),
        Inches(0.22),
        "List of Figures (--lof)",
        font_size=10,
        bold=True,
        color=ACCENT2,
    )
    sy += Inches(0.25)
    lof_entries = [
        ("Figure 1", "Block diagram", "12"),
        ("Figure 2", "Tile partitioning", "18"),
        ("Figure 3", "Transform tree", "24"),
    ]
    for fig, title, page in lof_entries:
        add_text_box(
            slide,
            px + Inches(0.3),
            sy,
            Inches(0.8),
            Inches(0.18),
            fig,
            font_size=9,
            bold=True,
            color=TEXT,
        )
        add_text_box(
            slide,
            px + Inches(1.1),
            sy,
            Inches(1.4),
            Inches(0.18),
            title,
            font_size=9,
            color=TEXT_LIGHT,
        )
        add_text_box(
            slide,
            px + Inches(3.0),
            sy,
            Inches(0.4),
            Inches(0.18),
            page,
            font_size=9,
            color=TEXT,
            alignment=PP_ALIGN.RIGHT,
        )
        sy += Inches(0.2)

    sy += Inches(0.08)
    add_text_box(
        slide,
        px + Inches(0.2),
        sy,
        Inches(3.3),
        Inches(0.22),
        "List of Tables (--lot)",
        font_size=10,
        bold=True,
        color=ACCENT3,
    )
    sy += Inches(0.25)
    lot_entries = [
        ("Table 1", "Syntax elements", "8"),
        ("Table 2", "Color primaries", "15"),
        ("Table 3", "Profile levels", "30"),
    ]
    for tab, title, page in lot_entries:
        add_text_box(
            slide,
            px + Inches(0.3),
            sy,
            Inches(0.8),
            Inches(0.18),
            tab,
            font_size=9,
            bold=True,
            color=TEXT,
        )
        add_text_box(
            slide,
            px + Inches(1.1),
            sy,
            Inches(1.4),
            Inches(0.18),
            title,
            font_size=9,
            color=TEXT_LIGHT,
        )
        add_text_box(
            slide,
            px + Inches(3.0),
            sy,
            Inches(0.4),
            Inches(0.18),
            page,
            font_size=9,
            color=TEXT,
            alignment=PP_ALIGN.RIGHT,
        )
        sy += Inches(0.2)


def slide_12_workflow(prs, layout):
    """Build Workflow Features."""
    slide = prs.slides.add_slide(layout)
    slide_title_bar(
        slide,
        "Build Workflow Features",
        "Incremental builds, watch mode, multipage, section filtering, and more",
    )

    workflow_features = [
        (
            "Incremental Builds",
            ACCENT,
            [
                "--incremental",
                "Skip compilation if sources unchanged",
                "Compares file timestamps",
                "Fast edit-preview cycles",
            ],
        ),
        (
            "Watch Mode",
            ACCENT2,
            [
                "--watch [--watch-interval N]",
                "Auto-rebuild on .bs file changes",
                "Configurable polling interval",
                "Combines with --incremental",
            ],
        ),
        (
            "Parallel Outputs",
            ACCENT3,
            [
                "--parallel-outputs",
                "Run PDF, DOCX, standalone concurrently",
                "ThreadPoolExecutor for I/O tasks",
                "Major speedup on multi-output builds",
            ],
        ),
        (
            "Multipage HTML",
            ACCENT4,
            [
                "--multipage",
                "Per-section files with navigation",
                "Sidebar TOC (left/right)",
                "Breadcrumb trail, prev/next links",
            ],
        ),
        (
            "Section Filtering",
            ACCENT5,
            [
                "--include-sections PATTERNS",
                "--exclude-sections PATTERNS",
                "Shell-style glob support",
                "Build subset of specification",
            ],
        ),
        (
            "Build Reports",
            TEAL,
            [
                "--build-report html|json|both",
                "Section stats, quality results",
                "Timing breakdown, messages",
                "CLI flags for reproducibility",
            ],
        ),
    ]

    card_w = Inches(3.8)
    card_h = Inches(2.15)

    for i, (title, color, items) in enumerate(workflow_features):
        col = i % 3
        row = i // 3
        x = Inches(0.55) + col * (card_w + Inches(0.3))
        y = Inches(1.6) + row * (card_h + Inches(0.2))
        add_card(
            slide,
            x,
            y,
            card_w,
            card_h,
            title,
            items,
            accent_color=color,
            title_size=14,
            body_size=11,
        )


def slide_13_source(prs, layout):
    """Source Processing & SDL."""
    slide = prs.slides.add_slide(layout)
    slide_title_bar(
        slide,
        "Source Processing & SDL",
        "From Bikeshed .bs files to compiled HTML with SDL syntax tables",
    )

    steps = [
        ("1", "Manifest", "Chapter ordering\nvia manifest.txt", ACCENT),
        ("2", "Merge", "Combine .bs files\ninto index.bs", ACCENT),
        ("3", "SDL Convert", "Descriptors to\nformatted tables", ACCENT2),
        ("4", "Bikeshed", "Compile to HTML\nwith cross-refs", ACCENT2),
        ("5", "Post-Process", "Enhance and\nvalidate the DOM", ACCENT3),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        x = Inches(0.4) + i * Inches(2.5)
        y = Inches(1.7)
        add_icon_circle(slide, x + Inches(0.75), y, Inches(0.55), color, num, label_size=22)
        add_text_box(
            slide,
            x,
            y + Inches(0.7),
            Inches(2.1),
            Inches(0.3),
            title,
            font_size=14,
            bold=True,
            color=color,
            alignment=PP_ALIGN.CENTER,
        )
        add_text_box(
            slide,
            x,
            y + Inches(1.05),
            Inches(2.1),
            Inches(0.6),
            desc,
            font_size=11,
            color=TEXT_LIGHT,
            alignment=PP_ALIGN.CENTER,
        )
        if i < len(steps) - 1:
            add_shape(
                slide,
                MSO_SHAPE.RIGHT_ARROW,
                x + Inches(2.15),
                y + Inches(0.12),
                Inches(0.28),
                Inches(0.3),
                fill_color=BORDER,
            )

    add_text_box(
        slide,
        Inches(0.5),
        Inches(3.7),
        Inches(6),
        Inches(0.35),
        "SOURCE & PROCESSING OPTIONS",
        font_size=13,
        bold=True,
        color=ACCENT,
    )

    src_opts = [
        ("--sdl", "Disable SDL table conversion (enabled by default)"),
        ("--compact", "Extract Section 9 tables to .h files"),
        ("--remove_editor_notes", "Strip editor notes from output"),
        ("--striped_code_blocks", "Alternating-color code block rows"),
        ("--auto-indent-code", "Auto-indent C/C++ by brace depth"),
        ("--manifest / --no-manifest", "Control chapter ordering"),
        ("--bikeshed-die-on LEVEL", "Set Bikeshed error threshold"),
        ("--externalize_resources", "Separate CSS/JS files"),
        ("--minify", "Minify HTML output"),
        ("--mobile_optimized", "Collapsible sections for mobile"),
    ]

    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(4.1),
        Inches(12),
        Inches(0.35),
        fill_color=NEAR_BLACK,
    )
    add_text_box(
        slide,
        Inches(0.6),
        Inches(4.13),
        Inches(3),
        Inches(0.3),
        "Flag",
        font_size=10,
        bold=True,
        color=WHITE,
        font_name="Consolas",
    )
    add_text_box(
        slide,
        Inches(3.8),
        Inches(4.13),
        Inches(3),
        Inches(0.3),
        "Description",
        font_size=10,
        bold=True,
        color=WHITE,
    )
    add_text_box(
        slide,
        Inches(6.9),
        Inches(4.13),
        Inches(3),
        Inches(0.3),
        "Flag",
        font_size=10,
        bold=True,
        color=WHITE,
        font_name="Consolas",
    )
    add_text_box(
        slide,
        Inches(10.1),
        Inches(4.13),
        Inches(3),
        Inches(0.3),
        "Description",
        font_size=10,
        bold=True,
        color=WHITE,
    )

    for i, (flag, desc) in enumerate(src_opts):
        col = i // 5
        row = i % 5
        iy = Inches(4.5) + row * Inches(0.42)
        bg = CARD_BG if row % 2 == 0 else WHITE
        fx = Inches(0.5) + col * Inches(6.4)
        add_shape(slide, MSO_SHAPE.RECTANGLE, fx, iy, Inches(6.3), Inches(0.42), fill_color=bg)
        add_text_box(
            slide,
            fx + Inches(0.1),
            iy + Inches(0.07),
            Inches(3.1),
            Inches(0.3),
            flag,
            font_size=10,
            color=TEXT,
            font_name="Consolas",
        )
        add_text_box(
            slide,
            fx + Inches(3.3),
            iy + Inches(0.07),
            Inches(2.8),
            Inches(0.3),
            desc,
            font_size=10,
            color=TEXT_LIGHT,
        )


def slide_14_build_report(prs, layout):
    """Build Report Dashboard."""
    slide = prs.slides.add_slide(layout)
    slide_title_bar(
        slide,
        "Build Report Dashboard",
        "--build-report html|json|both  \u2022  Comprehensive build summary for review and CI",
    )

    dx, dy = Inches(0.4), Inches(1.5)
    dw = Inches(12.5)
    _, dc_y = mockup_frame(
        slide, dx, dy, dw, Inches(5.5), "Build Report  \u2014  build_report.html"
    )

    add_text_box(
        slide,
        dx + Inches(0.2),
        dc_y + Inches(0.1),
        Inches(4),
        Inches(0.3),
        "Build Report",
        font_size=16,
        bold=True,
        color=NEAR_BLACK,
    )

    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        dx + Inches(10.5),
        dc_y + Inches(0.1),
        Inches(1.5),
        Inches(0.35),
        fill_color=CB_YELLOW,
    )
    add_text_box(
        slide,
        dx + Inches(10.5),
        dc_y + Inches(0.1),
        Inches(1.5),
        Inches(0.35),
        "WARNINGS",
        font_size=11,
        bold=True,
        color=CB_YELLOW_DARK,
        alignment=PP_ALIGN.CENTER,
    )

    add_text_box(
        slide,
        dx + Inches(0.2),
        dc_y + Inches(0.42),
        Inches(8),
        Inches(0.2),
        "Bicycle Specification  \u2014  main@fae537e  \u2022  2026-03-15  \u2022  4.2s total",
        font_size=10,
        color=MUTED,
    )

    cards_data = [
        ("142", "Sections", NEAR_BLACK, CARD_BG),
        ("48,350", "Words", NEAR_BLACK, CARD_BG),
        ("6", "Enhancements", NEAR_BLACK, CARD_BG),
        ("3", "Broken Refs", REPORT_WARN, RGBColor(0xFE, 0xFC, 0xBF)),
        ("0", "Broken Images", REPORT_OK, CARD_BG),
        ("5", "SDL Issues", REPORT_WARN, RGBColor(0xFE, 0xFC, 0xBF)),
        ("12", "Warnings", REPORT_WARN, RGBColor(0xFE, 0xFC, 0xBF)),
        ("0", "Errors", REPORT_OK, CARD_BG),
    ]
    card_w_small = Inches(1.4)
    card_h_small = Inches(0.85)
    cards_y = dc_y + Inches(0.7)
    for i, (value, label, val_color, bg_color) in enumerate(cards_data):
        cx = dx + Inches(0.2) + i * (card_w_small + Inches(0.12))
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            cx,
            cards_y,
            card_w_small,
            card_h_small,
            fill_color=bg_color,
            line_color=BORDER,
            line_width=Pt(1),
        )
        add_text_box(
            slide,
            cx,
            cards_y + Inches(0.08),
            card_w_small,
            Inches(0.4),
            value,
            font_size=22,
            bold=True,
            color=val_color,
            alignment=PP_ALIGN.CENTER,
        )
        add_text_box(
            slide,
            cx,
            cards_y + Inches(0.5),
            card_w_small,
            Inches(0.25),
            label,
            font_size=10,
            color=MUTED,
            alignment=PP_ALIGN.CENTER,
        )

    # Timing
    timing_y = cards_y + card_h_small + Inches(0.15)
    add_text_box(
        slide,
        dx + Inches(0.2),
        timing_y,
        Inches(3),
        Inches(0.25),
        "Timing",
        font_size=13,
        bold=True,
        color=NEAR_BLACK,
    )

    timing_data = [
        ("Merge .bs files", 0.12, 3),
        ("Bikeshed compile", 2.45, 58),
        ("Quality checks", 0.85, 20),
        ("Enhancements", 0.42, 10),
        ("PDF generation", 0.38, 9),
    ]
    ty = timing_y + Inches(0.3)
    for step, secs, pct in timing_data:
        add_text_box(
            slide, dx + Inches(0.3), ty, Inches(1.8), Inches(0.24), step, font_size=10, color=TEXT
        )
        add_text_box(
            slide,
            dx + Inches(2.2),
            ty,
            Inches(0.6),
            Inches(0.24),
            f"{secs:.2f}s",
            font_size=10,
            color=TEXT_LIGHT,
            alignment=PP_ALIGN.RIGHT,
        )
        add_text_box(
            slide,
            dx + Inches(2.9),
            ty,
            Inches(0.5),
            Inches(0.24),
            f"{pct}%",
            font_size=10,
            color=MUTED,
            alignment=PP_ALIGN.RIGHT,
        )
        bar_max = Inches(1.8)
        bar_w = max(Pt(3), int(bar_max * pct / 100))
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            dx + Inches(3.5),
            ty + Inches(0.03),
            bar_w,
            Inches(0.18),
            fill_color=ACCENT,
        )
        ty += Inches(0.28)

    # Broken refs
    qc_x = dx + Inches(5.8)
    add_text_box(
        slide,
        qc_x,
        timing_y,
        Inches(3),
        Inches(0.25),
        "Broken Cross-References (3)",
        font_size=12,
        bold=True,
        color=REPORT_WARN,
    )

    qc_header_y = timing_y + Inches(0.3)
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        qc_x,
        qc_header_y,
        Inches(6.3),
        Inches(0.3),
        fill_color=NEAR_BLACK,
    )
    for cx2, ct in [
        (qc_x + Inches(0.1), "href"),
        (qc_x + Inches(2.2), "Link Text"),
        (qc_x + Inches(4.3), "Near Section"),
    ]:
        add_text_box(
            slide,
            cx2,
            qc_header_y + Inches(0.03),
            Inches(2.0),
            Inches(0.26),
            ct,
            font_size=10,
            bold=True,
            color=WHITE,
        )

    ref_rows = [
        ("#section-tile-info", "tile_info()", "\u00a76.5 Tile Decoding"),
        ("#table-loop-filter-params", "Table 12", "\u00a76.3 Loop Filter"),
        ("#fig-transform-tree", "Figure 7", "\u00a76.8 Transform"),
    ]
    rry2 = qc_header_y + Inches(0.32)
    for i, (href, link_text, section) in enumerate(ref_rows):
        bg = CARD_BG if i % 2 == 0 else WHITE
        add_shape(slide, MSO_SHAPE.RECTANGLE, qc_x, rry2, Inches(6.3), Inches(0.3), fill_color=bg)
        add_text_box(
            slide,
            qc_x + Inches(0.1),
            rry2 + Inches(0.03),
            Inches(2.0),
            Inches(0.26),
            href,
            font_size=10,
            color=ACCENT5,
            font_name="Consolas",
        )
        add_text_box(
            slide,
            qc_x + Inches(2.2),
            rry2 + Inches(0.03),
            Inches(2.0),
            Inches(0.26),
            link_text,
            font_size=10,
            color=TEXT,
        )
        add_text_box(
            slide,
            qc_x + Inches(4.3),
            rry2 + Inches(0.03),
            Inches(2.0),
            Inches(0.26),
            section,
            font_size=10,
            color=TEXT_LIGHT,
        )
        rry2 += Inches(0.3)

    # SDL Issues
    rry2 += Inches(0.15)
    add_text_box(
        slide,
        qc_x,
        rry2,
        Inches(3),
        Inches(0.25),
        "SDL Reference Issues (5)",
        font_size=12,
        bold=True,
        color=REPORT_WARN,
    )
    rry2 += Inches(0.3)

    add_shape(
        slide, MSO_SHAPE.RECTANGLE, qc_x, rry2, Inches(6.3), Inches(0.3), fill_color=NEAR_BLACK
    )
    for cx2, ct in [
        (qc_x + Inches(0.1), "Unresolved Function"),
        (qc_x + Inches(2.5), "Referenced In"),
        (qc_x + Inches(4.5), "Section"),
    ]:
        add_text_box(
            slide,
            cx2,
            rry2 + Inches(0.03),
            Inches(2.0),
            Inches(0.26),
            ct,
            font_size=10,
            bold=True,
            color=WHITE,
        )

    sdl_rows = [
        ("get_position()", "frame_header", "\u00a76.2 Frame Header"),
        ("apply_grain()", "film_grain", "\u00a77.4 Film Grain"),
    ]
    rry2 += Inches(0.3)
    for i, (func, table, section) in enumerate(sdl_rows):
        bg = CARD_BG if i % 2 == 0 else WHITE
        add_shape(slide, MSO_SHAPE.RECTANGLE, qc_x, rry2, Inches(6.3), Inches(0.3), fill_color=bg)
        add_text_box(
            slide,
            qc_x + Inches(0.1),
            rry2 + Inches(0.03),
            Inches(2.2),
            Inches(0.26),
            func,
            font_size=10,
            color=ACCENT5,
            font_name="Consolas",
        )
        add_text_box(
            slide,
            qc_x + Inches(2.5),
            rry2 + Inches(0.03),
            Inches(1.8),
            Inches(0.26),
            table,
            font_size=10,
            color=TEXT,
            font_name="Consolas",
        )
        add_text_box(
            slide,
            qc_x + Inches(4.5),
            rry2 + Inches(0.03),
            Inches(1.8),
            Inches(0.26),
            section,
            font_size=10,
            color=TEXT_LIGHT,
        )
        rry2 += Inches(0.3)


def slide_15_getting_started(prs, layout):
    """Getting Started."""
    slide = prs.slides.add_slide(layout)
    slide_title_bar(slide, "Getting Started", "Installation, usage, and technology stack")

    add_text_box(
        slide,
        Inches(0.5),
        Inches(1.5),
        Inches(5),
        Inches(0.35),
        "INSTALLATION",
        font_size=14,
        bold=True,
        color=ACCENT,
    )

    install_cmds = [
        "pip install -e .             # Core",
        "pip install -e '.[pdf]'      # + PDF support",
        "pip install -e '.[test]'     # + Test deps",
        "pip install -e '.[dev]'      # Full dev setup",
    ]
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(1.9),
        Inches(5.5),
        Inches(1.6),
        fill_color=CODE_BG,
        line_color=BORDER,
        line_width=Pt(1),
    )
    y = Inches(2.05)
    for cmd in install_cmds:
        add_text_box(
            slide,
            Inches(0.7),
            y,
            Inches(5),
            Inches(0.3),
            cmd,
            font_size=11,
            color=ACCENT2,
            font_name="Consolas",
        )
        y += Inches(0.33)

    add_text_box(
        slide,
        Inches(7.0),
        Inches(1.5),
        Inches(5),
        Inches(0.35),
        "COMMON WORKFLOWS",
        font_size=14,
        bold=True,
        color=ACCENT,
    )

    examples = [
        ("Basic build:", "python compile.py"),
        ("Quick preview:", "python compile.py --profile quick"),
        ("Draft + PDF:", "python compile.py --profile draft --pdf"),
        ("Full review:", "python compile.py --profile review --build-report"),
        ("Publication:", "python compile.py --profile publication"),
        ("CI pipeline:", "python compile.py --all-checks --validate-refs-strict"),
        ("Watch mode:", "python compile.py --watch --incremental"),
    ]

    y = Inches(1.9)
    for label, cmd in examples:
        add_text_box(
            slide,
            Inches(7.0),
            y,
            Inches(2.0),
            Inches(0.25),
            label,
            font_size=11,
            bold=True,
            color=TEXT,
        )
        add_text_box(
            slide,
            Inches(7.0),
            y + Inches(0.25),
            Inches(5.5),
            Inches(0.25),
            cmd,
            font_size=10,
            color=ACCENT2,
            font_name="Consolas",
        )
        y += Inches(0.55)

    add_text_box(
        slide,
        Inches(0.5),
        Inches(3.9),
        Inches(5),
        Inches(0.35),
        "TECHNOLOGY STACK",
        font_size=14,
        bold=True,
        color=ACCENT,
    )

    stack_items = [
        ("Python 3.10+", "Core language", ACCENT),
        ("Bikeshed", "Spec preprocessor", ACCENT),
        ("BeautifulSoup4", "DOM manipulation", ACCENT2),
        ("lxml", "Fast HTML parsing", ACCENT2),
        ("WeasyPrint", "PDF engine (optional)", ACCENT3),
        ("Pygments", "Syntax highlighting", ACCENT3),
    ]

    for i, (tech, desc, color) in enumerate(stack_items):
        col = i // 3
        row = i % 3
        x = Inches(0.5) + col * Inches(3.0)
        iy = Inches(4.35) + row * Inches(0.55)
        add_shape(
            slide,
            MSO_SHAPE.OVAL,
            x,
            iy + Inches(0.05),
            Inches(0.12),
            Inches(0.12),
            fill_color=color,
        )
        add_text_box(
            slide,
            x + Inches(0.2),
            iy,
            Inches(1.4),
            Inches(0.25),
            tech,
            font_size=11,
            bold=True,
            color=TEXT,
        )
        add_text_box(
            slide,
            x + Inches(1.6),
            iy,
            Inches(1.3),
            Inches(0.25),
            desc,
            font_size=11,
            color=TEXT_LIGHT,
        )

    add_text_box(
        slide,
        Inches(0.5),
        Inches(6.2),
        Inches(11),
        Inches(0.5),
        "CI: GitHub Actions  \u2022  Python 3.10 / 3.12 / 3.13  \u2022  "
        "Ubuntu / macOS / Windows  \u2022  ruff lint + 767 tests",
        font_size=12,
        color=MUTED,
    )


# ═════════════════════════════════════════════════════════════════════════════
# Main
# ═════════════════════════════════════════════════════════════════════════════


def main():
    output_path = (
        Path(sys.argv[1])
        if len(sys.argv) > 1
        else (Path(__file__).resolve().parent.parent / "specbuild_features.pptx")
    )

    prs = Presentation()
    prs.slide_width = SLD_W
    prs.slide_height = SLD_H
    layout = prs.slide_layouts[6]  # blank

    # Build slides in final presentation order
    slide_01_title(prs, layout)  #  1: Title
    slide_02_architecture(prs, layout)  #  2: Architecture Overview
    slide_03_quality_checks(prs, layout)  #  3: Quality Checks (13)
    slide_04_qc_details(prs, layout)  #  4: QC Reports
    slide_05_enhancements(prs, layout)  #  5: Enhancements (18)
    slide_06_tooltips(prs, layout)  #  6: Tooltips & Search
    slide_07_change_tracking(prs, layout)  #  7: Change Bars & Diff
    slide_08_output_tasks(prs, layout)  #  8: Output Tasks (20)
    slide_09_profiles(prs, layout)  #  9: Build Profiles
    slide_10_pdf(prs, layout)  # 10: PDF Generation
    slide_11_pdf_enhancements(prs, layout)  # 11: PDF Enhancements
    slide_12_workflow(prs, layout)  # 12: Build Workflow
    slide_13_source(prs, layout)  # 13: Source Processing
    slide_14_build_report(prs, layout)  # 14: Build Report
    slide_15_getting_started(prs, layout)  # 15: Getting Started

    prs.save(str(output_path))
    print(f"Saved: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
